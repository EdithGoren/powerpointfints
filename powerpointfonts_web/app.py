from flask import Flask, request, render_template, send_file, jsonify
import os
import uuid
import time
from pptx import Presentation
from werkzeug.utils import secure_filename
import shutil
import threading

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # הגבלת גודל קובץ ל-16MB
app.config['ALLOWED_EXTENSIONS'] = {'pptx'}

# יצירת תיקיות נדרשות אם לא קיימות
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

def allowed_file(filename):
    """בדיקה האם סיומת הקובץ מותרת"""
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']

def clean_old_files():
    """ניקוי קבצים ישנים מתיקיית ההעלאות (מעל שעה)"""
    now = time.time()
    for filename in os.listdir(app.config['UPLOAD_FOLDER']):
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        if os.path.isfile(file_path) and os.path.getmtime(file_path) < now - 3600:
            try:
                os.remove(file_path)
            except:
                pass

def change_presentation_fonts(input_file, output_file, new_font):
    """שינוי פונטים בקובץ מצגת PowerPoint"""
    # טעינת המצגת
    prs = Presentation(input_file)
    
    # עיבוד כל השקופיות
    for slide in prs.slides:
        # עיבוד כל הצורות בשקופית
        for shape in slide.shapes:
            # אם הצורה היא טקסט
            if hasattr(shape, "text_frame") and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        # שינוי שם הפונט אך שמירה על גודל ופורמט
                        run.font.name = new_font
            
            # עיבוד טבלאות אם קיימות
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                # שינוי שם הפונט אך שמירה על גודל ופורמט
                                run.font.name = new_font
    
    # שמירת המצגת
    prs.save(output_file)

@app.route('/')
def index():
    """עמוד הבית של האפליקציה"""
    # הפעלת ניקוי קבצים ישנים ברקע
    cleaning_thread = threading.Thread(target=clean_old_files)
    cleaning_thread.start()
    
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload_file():
    """טיפול בהעלאת קובץ"""
    if 'file' not in request.files:
        return jsonify({'error': 'לא נמצא קובץ'}), 400
    
    file = request.files['file']
    
    if file.filename == '':
        return jsonify({'error': 'לא נבחר קובץ'}), 400
    
    if not allowed_file(file.filename):
        return jsonify({'error': 'סוג קובץ לא מורשה. אנא העלה קובץ PowerPoint (.pptx)'}), 400
    
    # יצירת שם קובץ ייחודי עם UUID
    unique_id = str(uuid.uuid4())
    original_filename = secure_filename(file.filename)
    filename_base, filename_ext = os.path.splitext(original_filename)
    
    # שמירת הקובץ המקורי
    input_filename = f"{unique_id}_input{filename_ext}"
    input_path = os.path.join(app.config['UPLOAD_FOLDER'], input_filename)
    file.save(input_path)
    
    # העברת פרטי הקובץ לעמוד הבא
    return jsonify({
        'success': True, 
        'file_id': unique_id, 
        'original_name': filename_base
    })

@app.route('/process', methods=['POST'])
def process_file():
    """עיבוד הקובץ ושינוי הפונטים"""
    data = request.json
    
    if not data or not all(key in data for key in ['file_id', 'font', 'original_name']):
        return jsonify({'error': 'נתונים חסרים'}), 400
    
    file_id = data['file_id']
    new_font = data['font']
    original_name = data['original_name']
    
    # בדיקה שהקובץ קיים
    input_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{file_id}_input.pptx")
    if not os.path.exists(input_path):
        return jsonify({'error': 'הקובץ לא נמצא'}), 404
    
    # יצירת שם קובץ לפלט
    output_filename = f"{file_id}_output.pptx"
    output_path = os.path.join(app.config['UPLOAD_FOLDER'], output_filename)
    
    try:
        # עיבוד הקובץ
        change_presentation_fonts(input_path, output_path, new_font)
        
        return jsonify({
            'success': True,
            'file_id': file_id,
            'original_name': original_name
        })
    
    except Exception as e:
        return jsonify({'error': f'שגיאה בעיבוד הקובץ: {str(e)}'}), 500

@app.route('/download/<file_id>/<original_name>')
def download_file(file_id, original_name):
    """הורדת הקובץ המעובד"""
    output_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{file_id}_output.pptx")
    
    if not os.path.exists(output_path):
        return "הקובץ לא נמצא", 404
    
    # הגדרת שם הקובץ להורדה
    download_name = f"{original_name}_new_font.pptx"
    
    return send_file(
        output_path, 
        as_attachment=True,
        download_name=download_name,
        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
    )

@app.route('/fonts')
def get_available_fonts():
    """החזרת רשימת פונטים נפוצים"""
    common_fonts = [
        "Arial", "Times New Roman", "Calibri", "Verdana", "Georgia", 
        "Tahoma", "Trebuchet MS", "Arial Black", "Impact", "Comic Sans MS",
        "Courier New", "Helvetica", "Palatino", "Garamond", "Book Antiqua",
        "Segoe UI", "Candara", "Consolas", "Constantia", "Corbel",
        "David", "Miriam", "Guttman Yad", "Levenim MT", "Narkisim",
        "Rod", "FrankRuehl", "Guttman Hatzvi"
    ]
    return jsonify(common_fonts)

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000) 