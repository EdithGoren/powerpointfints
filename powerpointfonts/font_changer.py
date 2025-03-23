import os
import sys
from pptx import Presentation
import tkinter as tk
from tkinter import filedialog, ttk, messagebox
from tkinter.font import families

class PowerPointFontChanger:
    def __init__(self, root):
        self.root = root
        self.root.title("PowerPoint Font Changer")
        self.root.geometry("600x400")
        
        # Set up the UI
        self.setup_ui()
        
    def setup_ui(self):
        # Frame for file selection
        file_frame = ttk.LabelFrame(self.root, text="Select PowerPoint File")
        file_frame.pack(fill="x", expand="yes", padx=20, pady=10)
        
        self.file_path = tk.StringVar()
        
        ttk.Label(file_frame, text="PowerPoint File:").grid(column=0, row=0, padx=5, pady=5, sticky=tk.W)
        ttk.Entry(file_frame, textvariable=self.file_path, width=50).grid(column=1, row=0, padx=5, pady=5)
        ttk.Button(file_frame, text="Browse...", command=self.browse_file).grid(column=2, row=0, padx=5, pady=5)
        
        # Frame for font selection
        font_frame = ttk.LabelFrame(self.root, text="Select New Font")
        font_frame.pack(fill="x", expand="yes", padx=20, pady=10)
        
        # Get available fonts
        self.available_fonts = sorted(list(families()))
        
        ttk.Label(font_frame, text="New Font:").grid(column=0, row=0, padx=5, pady=5, sticky=tk.W)
        self.font_combo = ttk.Combobox(font_frame, values=self.available_fonts, width=30)
        self.font_combo.grid(column=1, row=0, padx=5, pady=5)
        if self.available_fonts:
            self.font_combo.current(0)
        
        # Sample text to preview font
        ttk.Label(font_frame, text="Sample:").grid(column=0, row=1, padx=5, pady=5, sticky=tk.W)
        self.sample_text = tk.Text(font_frame, height=3, width=40)
        self.sample_text.grid(column=1, row=1, padx=5, pady=5)
        self.sample_text.insert(tk.END, "Sample Text אבג 123")
        
        ttk.Button(font_frame, text="Preview Font", command=self.preview_font).grid(column=2, row=1, padx=5, pady=5)
        
        # Process button
        ttk.Button(self.root, text="Change Fonts", command=self.change_fonts).pack(pady=20)
        
        # Status bar
        self.status_var = tk.StringVar()
        self.status_var.set("Ready")
        ttk.Label(self.root, textvariable=self.status_var, relief=tk.SUNKEN, anchor=tk.W).pack(side=tk.BOTTOM, fill=tk.X)
    
    def browse_file(self):
        filename = filedialog.askopenfilename(
            title="Select PowerPoint File",
            filetypes=(("PowerPoint files", "*.pptx"), ("All files", "*.*"))
        )
        if filename:
            self.file_path.set(filename)
    
    def preview_font(self):
        selected_font = self.font_combo.get()
        if selected_font:
            self.sample_text.configure(font=(selected_font, 10))
    
    def change_fonts(self):
        filepath = self.file_path.get()
        new_font = self.font_combo.get()
        
        if not filepath:
            messagebox.showerror("Error", "Please select a PowerPoint file")
            return
        
        if not os.path.exists(filepath):
            messagebox.showerror("Error", "File not found")
            return
        
        if not new_font:
            messagebox.showerror("Error", "Please select a font")
            return
        
        try:
            self.status_var.set("Processing...")
            self.root.update()
            
            # Create the output filename
            filename, ext = os.path.splitext(filepath)
            output_file = f"{filename}_new_font{ext}"
            
            # Process the file
            self.change_presentation_fonts(filepath, output_file, new_font)
            
            self.status_var.set(f"Done! Saved as {output_file}")
            messagebox.showinfo("Success", f"Font changed to {new_font}.\nSaved as {output_file}")
        
        except Exception as e:
            self.status_var.set("Error occurred")
            messagebox.showerror("Error", str(e))
    
    def change_presentation_fonts(self, input_file, output_file, new_font):
        # Load the presentation
        prs = Presentation(input_file)
        
        # Process all slides
        for slide in prs.slides:
            # Process all shapes in the slide
            for shape in slide.shapes:
                # If shape has text frame
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            # Change font name but keep size and formatting
                            run.font.name = new_font
                
                # Process tables if present
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    # Change font name but keep size and formatting
                                    run.font.name = new_font
        
        # Save the presentation
        prs.save(output_file)

def main():
    root = tk.Tk()
    app = PowerPointFontChanger(root)
    root.mainloop()

if __name__ == "__main__":
    main() 